import os
from typing import List, Dict, Any

class ExtractedSection:
    def __init__(self, text: str, page_number: int = 1, section_title: str = "General"):
        self.text = text
        self.page_number = page_number
        self.section_title = section_title

class DocumentExtractor:
    @staticmethod
    def extract_text(file_path: str, file_type: str) -> List[ExtractedSection]:
        ext = file_type.lower().replace(".", "")
        if ext == "pdf":
            return DocumentExtractor._extract_pdf(file_path)
        elif ext == "docx":
            return DocumentExtractor._extract_docx(file_path)
        elif ext == "pptx":
            return DocumentExtractor._extract_pptx(file_path)
        elif ext in ["txt", "md"]:
            return DocumentExtractor._extract_txt_md(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

    @staticmethod
    def _extract_pdf(file_path: str) -> List[ExtractedSection]:
        sections = []
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text("text").strip()
                if text:
                    # Clean repeated headers/footers
                    lines = [line.strip() for line in text.split("\n") if line.strip()]
                    cleaned_text = "\n".join(lines)
                    sections.append(ExtractedSection(
                        text=cleaned_text,
                        page_number=page_num + 1,
                        section_title=f"Page {page_num + 1}"
                    ))
            doc.close()
        except ImportError:
            # Fallback if PyMuPDF not available
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                sections.append(ExtractedSection(text=content, page_number=1, section_title="Document Content"))
        return sections

    @staticmethod
    def _extract_docx(file_path: str) -> List[ExtractedSection]:
        sections = []
        try:
            import docx
            doc = docx.Document(file_path)
            current_section = "Introduction"
            current_text = []
            page_estimate = 1
            paragraph_count = 0

            for p in doc.paragraphs:
                text = p.text.strip()
                if not text:
                    continue
                paragraph_count += 1
                if paragraph_count % 15 == 0:
                    page_estimate += 1

                if p.style.name.startswith("Heading"):
                    if current_text:
                        sections.append(ExtractedSection(
                            text="\n".join(current_text),
                            page_number=page_estimate,
                            section_title=current_section
                        ))
                        current_text = []
                    current_section = text
                else:
                    current_text.append(text)

            if current_text:
                sections.append(ExtractedSection(
                    text="\n".join(current_text),
                    page_number=page_estimate,
                    section_title=current_section
                ))
        except Exception as e:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                sections.append(ExtractedSection(text=f.read(), page_number=1, section_title="Main Content"))
        return sections

    @staticmethod
    def _extract_pptx(file_path: str) -> List[ExtractedSection]:
        sections = []
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = []
                slide_title = f"Slide {idx + 1}"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if shape == slide.shapes[0] and len(text) < 100:
                            slide_title = text
                        else:
                            slide_text.append(text)
                if slide_text:
                    sections.append(ExtractedSection(
                        text="\n".join(slide_text),
                        page_number=idx + 1,
                        section_title=slide_title
                    ))
        except Exception:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                sections.append(ExtractedSection(text=f.read(), page_number=1, section_title="Presentation Content"))
        return sections

    @staticmethod
    def _extract_txt_md(file_path: str) -> List[ExtractedSection]:
        sections = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        lines = content.split("\n")
        current_section = "General"
        current_lines = []
        page = 1

        for line in lines:
            line_str = line.strip()
            if line_str.startswith("#"):
                if current_lines:
                    sections.append(ExtractedSection(
                        text="\n".join(current_lines),
                        page_number=page,
                        section_title=current_section
                    ))
                    current_lines = []
                    page += 1
                current_section = line_str.lstrip("#").strip()
            elif line_str:
                current_lines.append(line_str)

        if current_lines:
            sections.append(ExtractedSection(
                text="\n".join(current_lines),
                page_number=page,
                section_title=current_section
            ))

        return sections if sections else [ExtractedSection(text=content, page_number=1, section_title="Full Document")]
